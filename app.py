"""
CV Tailoring Agent v2 — Streamlit front end
============================================
Runs TODAY with zero backend (mock mode) and switches automatically to the
real backend when core/ modules exist. This file defines the backend contract.

RUN IT:
    pip install streamlit pyyaml
    streamlit run app.py
Optional: pip install python-docx   (enables real .docx export; else .md export)
Optional: pip install python-docx pypdf python-pptx   (enables CV upload/import below)

BACKEND CONTRACT — implement these in Phase 1/2 and mock mode disappears:
    core/ingest.py    load_inventory(path: str) -> dict            # {"entries":[...]}
                      normalize_jd(text: str) -> str
    core/tailor.py    extract_keywords(jd: str) -> list[dict]      # [{"term","priority"}]
                      match(keywords, inventory) -> tuple[dict, list]  # (hits: {entry_id:[kw]}, gaps:[kw])
                      rewrite(entries: list, keywords: list) -> list[dict]  # [{"text","evidence_id"}]
    core/validate.py  validate_bullets(bullets, inventory) -> list[dict]
                      # each: {"text","evidence_id","status": "ok"|"blocked","violation": str|None}
    core/report.py    fit_score(hits, gaps, keywords) -> int       # 0-100 deterministic
                      export_docx(bullets: list[str]) -> bytes

GUARDRAILS BAKED INTO THIS UI (do not remove):
    - Export stays disabled until every changed bullet is accepted/rejected.
    - Blocked bullets are displayed, never hidden.
    - Gap keywords are reported, never auto-filled.
    - Real data lives in data/inventory.yaml (gitignored); this file ships none.
"""

from __future__ import annotations

import html
import io
import re
from pathlib import Path

import streamlit as st
import yaml

# ---------------------------------------------------------------------------
# Backend resolution: real core/ modules if present, otherwise mock mode.
# ---------------------------------------------------------------------------
MOCK_MODE = False
try:
    from core.ingest import load_inventory, normalize_jd            # type: ignore
    from core.tailor import extract_keywords, match, rewrite        # type: ignore
    from core.validate import validate_bullets                      # type: ignore
    from core.report import export_docx, fit_score                  # type: ignore
except ImportError:
    MOCK_MODE = True

# ---------------------------------------------------------------------------
# Mock backend (deleted from the flow automatically once core/ exists).
# The mock validate is a real miniature of the guard so the demo is honest.
# ---------------------------------------------------------------------------
SAMPLE_INVENTORY = {
    "entries": [
        {
            "id": "exp-tcs-001",
            "claim": "Managed 3 cross-functional Agile teams across ServiceNow ITSM, AWS and Salesforce for a Tier-1 telecom",
            "evidence": "TCS, Tier-1 Telecom client Malaysia, 2023-present; teams of 8/6/5; ServiceNow, AWS, Salesforce",
            "tags": ["agile", "servicenow", "aws", "salesforce", "scrum master", "telecom", "stakeholder management"],
        },
        {
            "id": "exp-umg-002",
            "claim": "Led BA workstream for royalty-platform migration at Universal Music Group",
            "evidence": "Universal Music Group; business analysis; royalty platform migration; requirements workshops",
            "tags": ["business analysis", "media", "migration", "requirements"],
        },
        {
            "id": "skill-jira-001",
            "claim": "Administered Jira boards and sprint dashboards for delivery reporting",
            "evidence": "Jira administration at TCS and UMG; sprint dashboards; delivery reporting",
            "tags": ["jira", "reporting", "scrum", "agile delivery"],
        },
    ]
}


def _mock_load_inventory(path: str) -> dict:
    p = Path(path)
    if p.exists():
        return yaml.safe_load(p.read_text(encoding="utf-8")) or {"entries": []}
    return SAMPLE_INVENTORY


def _mock_normalize_jd(text: str) -> str:
    return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", text)).strip()[:15000]


def _mock_extract_keywords(jd: str) -> list[dict]:
    """Pretend-LLM: known tags found in the JD = keywords; plus unknown tech-ish
    capitalised terms as gaps. Deterministic on purpose."""
    low = jd.lower()
    kws: list[dict] = []
    seen: set[str] = set()
    for e in st.session_state.inventory["entries"]:
        for t in e["tags"]:
            if t in low and t not in seen:
                seen.add(t)
                kws.append({"term": t, "priority": "must"})
    for term in re.findall(r"\b[A-Z][A-Za-z+#]{3,}\b", jd):
        t = term.lower()
        if t not in seen and t not in {"senior", "years", "agile"} and len(kws) < 18:
            seen.add(t)
            kws.append({"term": t, "priority": "nice"})
    return kws


def _mock_match(keywords: list[dict], inventory: dict, threshold: int = 85) -> tuple[dict, list]:
    hits: dict[str, list] = {}
    gaps: list[dict] = []
    for kw in keywords:
        found = None
        for e in inventory["entries"]:
            hay = " ".join(e["tags"]) + " " + e["claim"].lower()
            if kw["term"] in hay:
                found = e["id"]
                break
        if found:
            hits.setdefault(found, []).append(kw)
        else:
            gaps.append(kw)
    return hits, gaps


def _mock_rewrite(entries: list, keywords: list) -> list[dict]:
    bullets = [
        {
            "text": f"Led delivery grounded in: {e['claim']}",
            "evidence_id": e["id"],
        }
        for e in entries
    ]
    if entries:  # seeded fabrication so the guard is visible in the demo
        bullets.append({"text": "Certified Terraform practitioner managing IaC pipelines", "evidence_id": entries[0]["id"]})
    return bullets


_VERB_ALLOW = {"led", "certified", "improved", "managed", "administered",
               "delivered", "drove", "built", "owned", "coordinated"}


def _mock_validate(bullets: list[dict], inventory: dict) -> list[dict]:
    by_id = {e["id"]: e for e in inventory["entries"]}
    out = []
    for b in bullets:
        e = by_id.get(b["evidence_id"])
        if e is None:
            out.append({**b, "status": "blocked", "violation": f"unknown evidence_id {b['evidence_id']}"})
            continue
        hay = (e["claim"] + " " + e["evidence"]).lower()
        bad = None
        # numbers/metrics must exist in the cited evidence
        for n in re.findall(r"\d[\d,.]*%?", b["text"]):
            if n.rstrip("%") not in hay:
                bad = n
                break
        # named skills/tools (capitalized tokens) must exist too, verbs excused
        if bad is None:
            for tok in re.findall(r"\b[A-Z][A-Za-z]{3,}\b", b["text"]):
                if tok.lower() not in hay and tok.lower() not in _VERB_ALLOW:
                    bad = tok
                    break
        if bad:
            out.append({**b, "status": "blocked", "violation": f'token "{bad}" not found in evidence {e["id"]}'})
        else:
            out.append({**b, "status": "ok", "violation": None})
    return out


def _mock_fit_score(hits: dict, gaps: list, keywords: list) -> int:
    w = lambda k: 3 if k["priority"] == "must" else 1  # noqa: E731
    total = sum(w(k) for k in keywords) or 1
    got = sum(w(k) for kws in hits.values() for k in kws)
    return round(100 * got / total)


def _mock_export_docx(bullets: list[str]) -> bytes:
    try:
        from docx import Document  # type: ignore

        doc = Document()
        doc.add_heading("Tailored CV bullets", level=1)
        for b in bullets:
            doc.add_paragraph(b, style="List Bullet")
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        return ("# Tailored CV bullets\n" + "\n".join(f"- {b}" for b in bullets)).encode()


if MOCK_MODE:
    load_inventory, normalize_jd = _mock_load_inventory, _mock_normalize_jd
    extract_keywords, match, rewrite = _mock_extract_keywords, _mock_match, _mock_rewrite
    validate_bullets, fit_score, export_docx = _mock_validate, _mock_fit_score, _mock_export_docx

# ---------------------------------------------------------------------------
# Schema check (mirrors pydantic model to come in core/ingest.py)
# ---------------------------------------------------------------------------
def _read_docx_text(data: bytes) -> str:
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_pptx_text(data: bytes) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    lines.append(line)
    return "\n".join(lines)


def extract_cv_text(filename: str, data: bytes) -> str:
    """Best-effort text extraction for uploaded CVs. Raises if the required
    optional library isn't installed or the format is unsupported."""
    ext = Path(filename).suffix.lower()
    readers = {".docx": _read_docx_text, ".pdf": _read_pdf_text, ".pptx": _read_pptx_text}
    reader = readers.get(ext)
    if reader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    try:
        return reader(data)
    except ImportError as e:
        raise RuntimeError(
            f"Missing library to parse {ext} files — pip install "
            f"{'python-docx' if ext == '.docx' else 'pypdf' if ext == '.pdf' else 'python-pptx'}"
        ) from e


def schema_errors(inv: dict) -> list[str]:
    errs, ids = [], set()
    for i, e in enumerate(inv.get("entries", [])):
        for f in ("id", "claim", "evidence", "tags"):
            if not e.get(f):
                errs.append(f"entry {i}: missing '{f}'")
        eid = e.get("id")
        if eid:
            if eid in ids:
                errs.append(f"duplicate id '{eid}'")
            ids.add(eid)
        if e.get("tags") and not isinstance(e["tags"], list):
            errs.append(f"entry {i}: tags must be a list")
    return errs


# ---------------------------------------------------------------------------
# App state
# ---------------------------------------------------------------------------
st.set_page_config(page_title="CV Tailoring Agent", page_icon="🛡️", layout="wide")
TEAL, RED = "#0E7C7B", "#B4423A"

# ---- Visual polish layer (design tokens). Structural redesign = Phase 2, see UI_REDESIGN_SPEC.md ----
_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Source+Sans+3:wght@400;600;700&display=swap');
html { font-size:120%; }
html, body, [data-testid="stAppViewContainer"] { font-family:'Source Sans 3',system-ui,sans-serif; }
[data-testid="stAppViewContainer"] { background:#F7F9FB; color:#1F2933; }
[data-testid="stHeader"] { background:transparent; }
[data-testid="stSidebar"] { background:#1F2933; }
[data-testid="stSidebar"] * { color:#E6EAEE !important; font-family:'Source Sans 3',system-ui,sans-serif; }
/* Streamlit's icon glyphs (sidebar collapse arrow, uploader icon) are ligature text in a
   dedicated icon font — the blanket font-family rule above broke them into literal words. */
[data-testid="stIconMaterial"], .material-icons, .material-symbols-rounded {
  font-family: 'Material Symbols Rounded', 'Material Icons' !important;
}
h1,h2,h3 { font-family:'Source Sans 3',system-ui,sans-serif; font-weight:700; color:#1F2933 !important; letter-spacing:-.01em; }
p, span, label, div { color:#1F2933; }
[data-testid="stMetric"] { background:#FFFFFF; border:1px solid #E3E8EE; border-radius:10px;
  padding:14px 16px; box-shadow:0 1px 3px rgba(16,42,67,.06); }
[data-testid="stMetricValue"] { color:#0E7C7B !important; font-weight:700; }
[data-testid="stMetricLabel"] { color:#5B6472 !important; }

/* Buttons — explicit palette, both states, so text is never dark-on-dark */
.stButton>button, .stDownloadButton>button {
  background:#0E7C7B !important; color:#FFFFFF !important;
  border:1px solid #0E7C7B !important; border-radius:8px !important; font-weight:600 !important;
}
.stButton>button:hover, .stDownloadButton>button:hover {
  background:#0A5C5B !important; color:#FFFFFF !important; border-color:#0A5C5B !important;
}
.stButton>button:disabled, .stDownloadButton>button:disabled {
  background:#E9EBF0 !important; color:#9AA0AC !important; border-color:#E9EBF0 !important;
}
/* Secondary / plain buttons (e.g. checkbox-adjacent, radio) keep readable text on light bg */
[data-testid="stRadio"] label, [data-testid="stCheckbox"] label { color:#1F2933 !important; }

/* Sidebar logo — plain icon button, overrides the teal button fill above (declared later = wins) */
.st-key-logo_home button {
  background:transparent !important; border:none !important;
  padding:8px 0 4px 0 !important; box-shadow:none !important;
}
.st-key-logo_home button p { font-size:4.5rem !important; line-height:1 !important; color:#E6EAEE !important; }
.st-key-logo_home button:hover p { color:#2DBDBA !important; }
.st-key-logo_home { margin-bottom:-4px; }

/* File uploader (sidebar) — dark dropzone matches the sidebar card, teal Browse button */
[data-testid="stFileUploaderDropzone"] {
  background:#2A3441 !important; border:1px dashed #4A5568 !important; border-radius:8px !important;
}
[data-testid="stFileUploaderDropzoneInstructions"] * { color:#E6EAEE !important; }
[data-testid="stFileUploaderDropzoneInstructions"] svg { fill:#E6EAEE !important; }
[data-testid="stFileUploaderDropzone"] button {
  background:#0E7C7B !important; color:#FFFFFF !important;
  border:1px solid #0E7C7B !important; border-radius:8px !important; font-weight:600 !important;
  width:auto !important; height:auto !important; min-width:fit-content !important;
  white-space:nowrap !important; padding:6px 16px !important;
}
[data-testid="stFileUploaderDropzone"] button:hover {
  background:#0A5C5B !important; border-color:#0A5C5B !important;
}
[data-testid="stFileUploaderFile"] { background:#2A3441 !important; border-radius:8px !important; }
[data-testid="stFileUploaderFile"] * { color:#E6EAEE !important; }

textarea, input[type="text"] { border-radius:8px !important; color:#1F2933 !important; background:#FFFFFF !important; }

.cva-chip{display:inline-block;background:#E3F2F1;color:#0A5C5B;border:1px solid #9AD5D4;
  border-radius:14px;padding:2px 12px;margin:0 6px 6px 0;font-size:.85rem;font-weight:600}
.cva-chip-miss{background:#F9ECEB;color:#B4423A;border-color:#E5B7B2}
.cva-card{background:#fff;border:1px solid #E3E8EE;border-left:4px solid #0E7C7B;border-radius:10px;
  padding:12px 16px;margin-bottom:6px;box-shadow:0 1px 3px rgba(16,42,67,.06);color:#1F2933}
.cva-card-blocked{border-left-color:#B4423A;background:#FBF3F2}
.cva-ev{font-family:Consolas,monospace;font-size:.72rem;color:#0A5C5B;background:#E3F2F1;
  border-radius:4px;padding:1px 8px;letter-spacing:.03em}
.cva-ev-err{color:#B4423A;background:#F9ECEB}
.cva-gapbox{background:#FBF3F2;border:1px solid #E5B7B2;border-radius:10px;padding:12px 16px;margin:8px 0;color:#1F2933}

/* Sidebar upload callout — overrides the blanket light-text rule below so it reads on the dark card */
.cva-sidebar-cta{background:rgba(14,124,123,.22);border:1px solid #0E7C7B;border-radius:8px;
  padding:10px 12px;margin:10px 0;font-size:.85rem;line-height:1.4}
.cva-sidebar-cta, .cva-sidebar-cta *{color:#E6EAEE !important}
</style>
"""
st.markdown(_CSS, unsafe_allow_html=True)

ss = st.session_state
ss.setdefault("inventory", load_inventory("data/inventory.yaml"))
ss.setdefault("run", None)          # {"keywords","hits","gaps","score"}
ss.setdefault("run_id", 0)          # nonce: isolates widget keys per run (stale-state fix)
ss.setdefault("ticked", set())      # entry ids selected for rewrite
ss.setdefault("bullets", [])        # validated bullets
ss.setdefault("decisions", {})      # bullet index -> "accepted"|"rejected"

if "_nav_next" in ss:  # deferred nav: must land before the radio widget below is instantiated
    ss["nav_radio"] = ss.pop("_nav_next")

with st.sidebar:
    if st.button("🛡️", key="logo_home", help="Home", type="tertiary"):
        st.session_state["nav_radio"] = "📇 Inventory"
        st.rerun()
    st.markdown("### CV Tailoring Agent")
    st.caption(f"local · evidence-only · v2 · {len(ss.inventory.get('entries', []))} entries")
    screen = st.radio("Navigate", ["📇 Inventory", "🎯 Tailor", "✅ Review & Export"],
                       key="nav_radio", label_visibility="collapsed")
    if MOCK_MODE:
        st.warning("MOCK MODE — core/ not found. Build Phase 1 to go live.", icon="⚠️")
    st.divider()

    st.markdown("**Import CV**")
    uploaded = st.file_uploader(
        "Upload .docx / .pdf / .pptx", type=["docx", "pdf", "pptx"], key="cv_upload"
    )
    if uploaded is not None:
        raw = uploaded.getvalue()
        upload_dir = Path("data/uploads")
        upload_dir.mkdir(parents=True, exist_ok=True)
        save_path = upload_dir / uploaded.name
        save_path.write_bytes(raw)

        try:
            cv_text = extract_cv_text(uploaded.name, raw)
        except Exception as e:
            st.error(str(e))
            cv_text = ""

        if cv_text.strip():
            st.markdown(
                f"<div class='cva-sidebar-cta'>📄 Extracted {len(cv_text):,} characters from "
                f"<b>{html.escape(uploaded.name)}</b> — add tags below, then add it as an entry.</div>",
                unsafe_allow_html=True,
            )
            tags_input = st.text_input(
                "Tags (comma-separated, required)", key="cv_add_tags", placeholder="e.g. scrum-master, agile"
            )
            tags = [t.strip() for t in tags_input.split(",") if t.strip()]
            if st.button("➕ Add as inventory entry", key="cv_add_entry",
                         use_container_width=True, disabled=not tags):
                existing_ids = {e["id"] for e in ss.inventory.get("entries", [])}
                base_id = f"upload-{Path(uploaded.name).stem.lower().replace(' ', '-')}"
                uid, n = base_id, 1
                while uid in existing_ids:
                    n += 1
                    uid = f"{base_id}-{n}"
                first_line = next((l.strip() for l in cv_text.splitlines() if l.strip()), uploaded.name)
                ss.inventory.setdefault("entries", []).append({
                    "id": uid,
                    "claim": first_line[:120],
                    "evidence": cv_text.strip()[:5000],
                    "tags": tags,
                })
                st.success(f"Added '{uid}' — refine claim/tags on the Inventory screen, then Save.")
            with st.expander("Preview extracted text"):
                st.text(cv_text[:3000] + ("…" if len(cv_text) > 3000 else ""))

errors = schema_errors(ss.inventory)

# ---------------------------------------------------------------------------
# S1 — INVENTORY
# ---------------------------------------------------------------------------
if screen == "📇 Inventory":
    st.subheader("Evidence inventory")
    st.caption("Everything a tailored CV is allowed to say lives here. No entry — no claim.")

    entries = ss.inventory.get("entries", [])
    if errors:
        st.error("SCHEMA INVALID — Tailor is disabled:\n\n- " + "\n- ".join(errors))
    elif entries:
        st.success("SCHEMA VALID ✓")

    ss.setdefault("show_manual_entry", False)
    if not entries and not ss.show_manual_entry:
        st.info("No entries yet — upload a CV from the sidebar, or add one manually below.")
        if st.button("➕ Add a blank entry"):
            ss.show_manual_entry = True
            st.rerun()
    else:
        scaffold = entries or [{"id": "", "claim": "", "evidence": "", "tags": []}]
        rows = [
            {"id": e["id"], "claim": e["claim"], "evidence": e["evidence"], "tags": ", ".join(e["tags"])}
            for e in scaffold
        ]
        edited = st.data_editor(rows, num_rows="dynamic", use_container_width=True, key="inv_editor")
        if any(r.get("id") for r in edited) and st.button("💾 Save your skills to compare against JD"):
            new_inv = {
                "entries": [
                    {"id": r["id"], "claim": r["claim"], "evidence": r["evidence"],
                     "tags": [t.strip() for t in str(r["tags"]).split(",") if t.strip()]}
                    for r in edited if r.get("id")
                ]
            }
            errs = schema_errors(new_inv)
            if errs:
                st.error("Not saved:\n\n- " + "\n- ".join(errs))
            else:
                Path("data").mkdir(exist_ok=True)
                Path("data/inventory.yaml").write_text(yaml.safe_dump(new_inv, sort_keys=False), encoding="utf-8")
                ss.inventory = new_inv
                ss.show_manual_entry = False
                ss.run, ss.bullets, ss.decisions, ss.ticked = None, [], {}, set()  # invalidate stale run
                st.success("Saved. Previous tailor run cleared.")

# ---------------------------------------------------------------------------
# S2 — TAILOR
# ---------------------------------------------------------------------------
elif screen == "🎯 Tailor":
    st.subheader("Tailor to a job description")
    if errors:
        st.error("Fix the inventory schema first (see Inventory screen).")
        st.stop()
    st.caption("The JD is untrusted input — treated as data, never as instructions.")
    jd = st.text_area("Paste job description", height=170, placeholder="Senior Scrum Master — ServiceNow & Terraform…")

    if st.button("Run tailor →", type="primary", disabled=not jd.strip()):
        try:
            with st.spinner("Extracting keywords…"):
                jd_n = normalize_jd(jd)
                keywords = extract_keywords(jd_n)
                hits, gaps = match(keywords, ss.inventory)
        except Exception as e:
            st.error(f"Couldn't run tailor: {e}")
        else:
            ss.run = {"keywords": keywords, "hits": hits, "gaps": gaps,
                      "score": fit_score(hits, gaps, keywords)}
            ss.run_id += 1
            ss.ticked = set(hits.keys())
            ss.bullets, ss.decisions = [], {}

    if ss.run:
        r = ss.run
        c1, c2, c3 = st.columns(3)
        c1.metric("Fit score (deterministic)", f"{r['score']}/100")
        c2.metric("Matched keywords", sum(len(v) for v in r["hits"].values()))
        c3.metric("Gaps — never filled", len(r["gaps"]))

        seen_terms: set[str] = set()
        chip_html = ""
        for kws_list in r["hits"].values():
            for k in kws_list:
                if k["term"] not in seen_terms:
                    seen_terms.add(k["term"])
                    chip_html += f"<span class='cva-chip'>{html.escape(k['term'])}</span>"
        if chip_html:
            st.markdown(chip_html, unsafe_allow_html=True)

        st.markdown("**Matched evidence** — untick to exclude, tick to force-include (your override):")
        by_id = {e["id"]: e for e in ss.inventory["entries"]}
        for e in ss.inventory["entries"]:
            kws = ", ".join(k["term"] for k in r["hits"].get(e["id"], []))
            label = f"`{e['id']}` — {e['claim'][:90]}" + (f"  ·  _{kws}_" if kws else "  ·  _no keyword match — tick only if truly relevant_")
            on = st.checkbox(label, value=e["id"] in ss.ticked, key=f"tick_{ss.run_id}_{e['id']}")
            if on:
                ss.ticked.add(e["id"])
            else:
                ss.ticked.discard(e["id"])

        if r["gaps"]:
            gap_chips = "".join(
                f"<span class='cva-chip cva-chip-miss'>{html.escape(g['term'])} · {g['priority']}</span>"
                for g in r["gaps"]
            )
            st.markdown(
                f"<div class='cva-gapbox'><b style='color:{RED}'>Gap report — no evidence for:</b><br>{gap_chips}"
                f"<div style='color:{RED};font-size:.85rem;font-style:italic'>The agent will not fabricate these — "
                "add evidence, reframe honestly, or accept the gap.</div></div>",
                unsafe_allow_html=True,
            )

        if st.button("Rewrite selected evidence →", disabled=not ss.ticked):
            entries = [by_id[i] for i in ss.ticked if i in by_id]  # guard: ids may have been deleted
            try:
                with st.spinner("Rewriting bullets…"):
                    raw = rewrite(entries, r["keywords"])
                    ss.bullets = validate_bullets(raw, ss.inventory)
            except Exception as e:
                st.error(f"Couldn't rewrite: {e}")
            else:
                ss.decisions = {}
                ss["_nav_next"] = "✅ Review & Export"
                st.rerun()

# ---------------------------------------------------------------------------
# S3 — REVIEW & EXPORT
# ---------------------------------------------------------------------------
else:
    st.subheader("Review every change")
    if not ss.bullets:
        st.info("No run yet — go to Tailor first.")
        st.stop()

    pending = 0
    for i, b in enumerate(ss.bullets):
        left, right = st.columns([1, 1])
        txt, vio, eid = html.escape(b["text"]), html.escape(b.get("violation") or ""), html.escape(b["evidence_id"])
        if b["status"] == "blocked":
            right.markdown(
                f"<div class='cva-card cva-card-blocked'>"
                f"<span class='cva-ev cva-ev-err'>BLOCKED BY VALIDATE.PY</span><br>"
                f"<s style='color:#808495'>{txt}</s><br>"
                f"<span style='color:{RED};font-size:.85rem'>{vio} — original kept, violation logged</span></div>",
                unsafe_allow_html=True,
            )
            left.caption("guard rejection — no action needed")
            continue
        right.markdown(
            f"<div class='cva-card'>"
            f"<span class='cva-ev'>evidence: {eid}</span><br>{txt}</div>",
            unsafe_allow_html=True,
        )
        choice = left.radio("Decision", ["pending", "accept", "reject"], key=f"dec_{ss.run_id}_{i}", horizontal=True)
        ss.decisions[i] = choice
        pending += choice == "pending"

    st.divider()
    accepted = [b["text"] for i, b in enumerate(ss.bullets) if ss.decisions.get(i) == "accept"]
    if pending:
        st.button(f"⬇ Export DOCX — locked ({pending} pending)", disabled=True)
        st.caption("Guardrail #2: nothing exports until every bullet has your explicit decision.")
    else:
        data = export_docx(accepted)
        is_docx = data[:2] == b"PK"
        st.download_button(
            "⬇ Export tailored bullets",
            data=data,
            file_name="tailored_cv." + ("docx" if is_docx else "md"),
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document" if is_docx else "text/markdown",
        )
        st.caption(f"{len(accepted)} accepted bullets.")
