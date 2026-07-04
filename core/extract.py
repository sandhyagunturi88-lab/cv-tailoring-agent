"""
core/extract.py — best-effort text extraction from uploaded CVs.

Returns plain text for the caller to turn into an inventory entry. This is
dumb, deterministic extraction: no LLM calls here.
"""

from __future__ import annotations

import io
from pathlib import Path


def _read_docx_text(data: bytes) -> str:
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_pptx_text(data: bytes) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    lines.append(line)
    return "\n".join(lines)


def extract_cv_text(filename: str, data: bytes) -> str:
    """Best-effort text extraction for uploaded CVs. Raises if the required
    optional library isn't installed or the format is unsupported."""
    ext = Path(filename).suffix.lower()
    readers = {".docx": _read_docx_text, ".pdf": _read_pdf_text, ".pptx": _read_pptx_text}
    reader = readers.get(ext)
    if reader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    try:
        return reader(data)
    except ImportError as e:
        raise RuntimeError(
            f"Missing library to parse {ext} files — pip install "
            f"{'python-docx' if ext == '.docx' else 'pypdf' if ext == '.pdf' else 'python-pptx'}"
        ) from e
